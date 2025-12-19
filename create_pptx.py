from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# --- 設計大師的調色盤 ---
BG_COLOR = RGBColor(18, 18, 18)        # 極致黑 (Premium Black)
ACCENT_GOLD = RGBColor(212, 175, 55)   # 香檳金 (Champagne Gold)
ACCENT_BLUE = RGBColor(52, 152, 219)   # 科技藍 (Tech Blue) - 用於少量點綴
TEXT_WHITE = RGBColor(240, 240, 240)   # 珍珠白 (Pearl White)
TEXT_GREY = RGBColor(160, 160, 160)    # 銀灰 (Silver Grey)
CARD_BG = RGBColor(35, 35, 35)         # 卡片深灰 (Card Grey)

def create_premium_deck():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # --- 核心設計元素函式 ---
    def set_slide_background(slide):
        """設定沉穩的深色背景"""
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = BG_COLOR

    def add_fluid_elements(slide):
        """加入具有設計感的流線背景元素 (Abstract Fluid Shapes)"""
        # 右上角的裝飾流線
        shape = slide.shapes.add_shape(MSO_SHAPE.CHORD, Inches(10), Inches(-1), Inches(4), Inches(4))
        shape.fill.solid()
        shape.fill.fore_color.rgb = ACCENT_GOLD
        shape.rotation = 45
        # 模擬半透明效果 (透過調整顏色亮度，因為 python-pptx 直接設 alpha 較複雜)
        # 這裡我們用一個較暗的金屬色來模擬
        shape.fill.fore_color.rgb = RGBColor(60, 50, 20) 
        shape.line.fill.background()

        # 左下角的呼應流線
        shape2 = slide.shapes.add_shape(MSO_SHAPE.MOON, Inches(-1), Inches(5), Inches(3), Inches(3))
        shape2.fill.solid()
        shape2.fill.fore_color.rgb = RGBColor(40, 40, 40)
        shape2.rotation = -30
        shape2.line.fill.background()
        
        # 頂部金色細線 (Header Line)
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.2), Inches(1.5), Inches(0.03))
        line.fill.solid()
        line.fill.fore_color.rgb = ACCENT_GOLD
        line.line.fill.background()

    def add_title(slide, text, subtitle=None):
        """加入高對比標題"""
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(10), Inches(1))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.name = 'Arial'
        p.font.color.rgb = TEXT_WHITE
        
        if subtitle:
            sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(10), Inches(0.5))
            sf = sub_box.text_frame
            p = sf.paragraphs[0]
            p.text = subtitle
            p.font.size = Pt(16)
            p.font.color.rgb = ACCENT_GOLD

    def add_card(slide, x, y, w, h, title, content):
        """加入懸浮卡片設計"""
        # 卡片背景
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
        shape.fill.solid()
        shape.fill.fore_color.rgb = CARD_BG
        shape.line.color.rgb = RGBColor(60, 60, 60)
        shape.line.width = Pt(0.5)
        
        # 卡片標題
        tb_title = slide.shapes.add_textbox(x + Inches(0.2), y + Inches(0.2), w - Inches(0.4), Inches(0.5))
        p = tb_title.text_frame.paragraphs[0]
        p.text = title
        p.font.bold = True
        p.font.size = Pt(14)
        p.font.color.rgb = ACCENT_GOLD
        
        # 卡片內容
        tb_content = slide.shapes.add_textbox(x + Inches(0.2), y + Inches(0.7), w - Inches(0.4), h - Inches(0.8))
        tf = tb_content.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = content
        p.font.size = Pt(11)
        p.font.color.rgb = TEXT_GREY

    # ================= 投影片製作 =================
    gap = Inches(0.5)

    # --- Slide 1: 封面 (大氣、留白) ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide)
    
    # 裝飾：巨大的流線圓弧
    big_circle = slide.shapes.add_shape(MSO_SHAPE.ARC, Inches(4), Inches(1), Inches(9), Inches(9))
    big_circle.line.color.rgb = ACCENT_GOLD
    big_circle.line.width = Pt(1.5)
    
    # 標題區
    title = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(2))
    p = title.text_frame.paragraphs[0]
    p.text = "CoreSelf AI"
    p.font.size = Pt(72)
    p.font.bold = True
    p.font.color.rgb = TEXT_WHITE
    
    sub = slide.shapes.add_textbox(Inches(1), Inches(3.8), Inches(8), Inches(1))
    p = sub.text_frame.paragraphs[0]
    p.text = "多角色心理增效 AI 顧問"
    p.font.size = Pt(28)
    p.font.color.rgb = ACCENT_GOLD
    
    tagline = slide.shapes.add_textbox(Inches(1), Inches(6), Inches(6), Inches(1))
    p = tagline.text_frame.paragraphs[0]
    p.text = "守住人生選擇的清醒\n減少角色內耗"
    p.font.size = Pt(14)
    p.font.color.rgb = TEXT_GREY

    # --- Slide 2: 核心定位 (Quote Style) ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide)
    add_fluid_elements(slide)
    add_title(slide, "核心定位", "OUR VALUE PROPOSITION")
    
    # 引用框
    quote_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.5), Inches(3), Inches(10.33), Inches(2))
    quote_bg.fill.background() # 透明
    quote_bg.line.color.rgb = ACCENT_GOLD
    quote_bg.line.width = Pt(2)
    
    quote = slide.shapes.add_textbox(Inches(2), Inches(3.5), Inches(9.33), Inches(1.5))
    p = quote.text_frame.paragraphs[0]
    p.text = "「不是交出人生給 AI，而是幫使用者在混亂時，\n守住選擇的清醒。」"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = TEXT_WHITE
    p.alignment = PP_ALIGN.CENTER

    # --- Slide 3: 核心痛點 (卡片式) ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide)
    add_fluid_elements(slide)
    add_title(slide, "核心痛點：現代人的多重角色困境", "THE PROBLEM")
    
    y_pos = Inches(2.5)
    w = Inches(3.5)
    h = Inches(3.5)
    
    add_card(slide, Inches(1), y_pos, w, h, "01. 多重角色衝突", 
             "主管 / 父母 / 子女 / 自我...\n角色切換頻繁，導致心理能量耗損，無法專注於當下。")
    add_card(slide, Inches(1)+w+gap, y_pos, w, h, "02. 決策疲勞", 
             "每日數百個微小決策耗盡意志力。\n現有工具（Notion/Calendar）只管時間，不管「心理能量」。")
    add_card(slide, Inches(1)+(w+gap)*2, y_pos, w, h, "03. 缺乏持續支持", 
             "心理諮詢成本高、頻率低。\n缺乏「即時反饋」與「數據化」的落地執行方案。")

    # --- Slide 4: 產品架構 (層疊視覺) ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide)
    add_fluid_elements(slide)
    add_title(slide, "產品架構：三層心理增效系統", "SOLUTION ARCHITECTURE")
    
    layer_w = Inches(7)
    layer_h = Inches(1.2)
    x_center = Inches(1.5)
    
    # 繪製層疊圖
    labels = [
        ("Interaction Flow (交互層)", "每日核心建議 • 衝突預警 • 決策回報"),
        ("Intelligence Layer (智能層)", "資訊整理 • 衝突對齊 • 心理權重計算"),
        ("User Context Layer (情境層)", "角色地圖 (Role Map) • 優先級 • 狀態標記")
    ]
    
    for i, (l_title, l_desc) in enumerate(labels):
        top = Inches(2.5) + (i * Inches(1.4))
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x_center, top, layer_w, layer_h)
        # 漸層模擬：越底層越暗
        shape.fill.solid()
        col_val = 60 - (i * 10)
        shape.fill.fore_color.rgb = RGBColor(col_val, col_val, col_val + 10) 
        
        # 文字
        tf = shape.text_frame
        p = tf.paragraphs[0]
        p.text = f"{l_title}\n{l_desc}"
        p.font.color.rgb = TEXT_WHITE if i == 0 else RGBColor(200,200,200)
        p.font.size = Pt(12)
        p.alignment = PP_ALIGN.CENTER
        
    # MVP 說明 (右側)
    add_card(slide, Inches(9), Inches(2.5), Inches(3.5), Inches(4), "MVP 核心功能", 
             "• 角色盤點器 (Role Map Lite)\n\n• 今日角色輪值建議\n\n• 衝突分析 Dashboard\n\n• 簡易決策回報")

    # --- Slide 5: 護城河 (TRR 模型) ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide)
    add_fluid_elements(slide)
    add_title(slide, "AI 核心護城河：TRR 模型", "COMPETITIVE ADVANTAGE")
    
    # 中心圓
    center_circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(5), Inches(2.5), Inches(3), Inches(3))
    center_circle.fill.solid()
    center_circle.fill.fore_color.rgb = RGBColor(0,0,0)
    center_circle.line.color.rgb = ACCENT_GOLD
    center_circle.line.width = Pt(3)
    
    tf = center_circle.text_frame
    p = tf.paragraphs[0]
    p.text = "TRR\nData\nLoop"
    p.font.bold = True
    p.font.size = Pt(24)
    p.font.color.rgb = ACCENT_GOLD
    p.alignment = PP_ALIGN.CENTER
    
    # 兩側說明
    add_card(slide, Inches(1), Inches(3), Inches(3.5), Inches(2), "心理學權重模型", 
             "量化人生意圖，將抽象價值觀轉化為每日權重建議。")
    add_card(slide, Inches(8.5), Inches(3), Inches(3.5), Inches(2), "多層安全防火牆", 
             "Level 3 人工審核 + 倫理導師 + 法律護欄。")

    # --- Slide 6: 商業模式 (高價值階梯) ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide)
    add_fluid_elements(slide)
    add_title(slide, "商業模式：階梯式價值主張", "BUSINESS MODEL")
    
    # 階梯視覺
    # Step 1
    s1 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(5.5), Inches(3), Inches(1.5))
    s1.fill.solid()
    s1.fill.fore_color.rgb = RGBColor(60, 60, 60)
    s1.text_frame.text = "入門試用 $10\n基本體驗"
    for p in s1.text_frame.paragraphs:
        p.font.color.rgb = TEXT_WHITE
        p.font.size = Pt(14)
        p.alignment = PP_ALIGN.CENTER
    
    # Step 2
    s2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4.2), Inches(4), Inches(3), Inches(3))
    s2.fill.solid()
    s2.fill.fore_color.rgb = RGBColor(80, 80, 80)
    s2.text_frame.text = "標準會員 $50\n習慣養成"
    for p in s2.text_frame.paragraphs:
        p.font.color.rgb = TEXT_WHITE
        p.font.size = Pt(14)
        p.alignment = PP_ALIGN.CENTER
    
    # Step 3
    s3 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(7.4), Inches(2.5), Inches(3.5), Inches(4.5))
    s3.fill.solid()
    s3.fill.fore_color.rgb = ACCENT_GOLD # Highlight
    tf = s3.text_frame
    p = tf.paragraphs[0]
    p.text = "高級會員\n$150 - $200"
    p.font.bold = True
    p.font.size = Pt(20)
    p.font.color.rgb = BG_COLOR # Dark text on gold
    p.alignment = PP_ALIGN.CENTER
    
    p2 = tf.add_paragraph()
    p2.text = "\n全 AI 建議 + 導師抽查\nLTV 核心來源"
    p2.font.size = Pt(12)
    p2.font.color.rgb = BG_COLOR
    p2.alignment = PP_ALIGN.CENTER

    # --- Slide 7: 獲客策略 ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide)
    add_fluid_elements(slide)
    add_title(slide, "獲客策略：規模化與自動化", "GO-TO-MARKET")
    
    col_w = Inches(3.5)
    add_card(slide, Inches(1), Inches(2.5), col_w, Inches(2.5), "1. 系統化推薦", 
             "利用高信任度用戶啟動飛輪。將推薦機制自動化，降低獲客成本。")
    add_card(slide, Inches(1)+col_w+gap, Inches(2.5), col_w, Inches(2.5), "2. 社群矩陣化", 
             "多平台內容經營。建立專業形象與用戶信任。")
    add_card(slide, Inches(1)+(col_w+gap)*2, Inches(2.5), col_w, Inches(2.5), "3. B2B 企業方案", 
             "企業員工心理健康方案。提供團體授權與數據分析。")

    # --- Slide 8: 團隊與願景 ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide)
    add_fluid_elements(slide)
    add_title(slide, "我們的願景", "OUR VISION")
    
    # 願景陳述
    vision_box = slide.shapes.add_textbox(Inches(1.5), Inches(3), Inches(10), Inches(2))
    tf = vision_box.text_frame
    p = tf.paragraphs[0]
    p.text = "讓每個人都能在人生的角色切換中，\n保持清醒、減少內耗，\n活出更有意義的生活。"
    p.font.size = Pt(28)
    p.font.color.rgb = TEXT_WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # 聯繫資訊
    contact = slide.shapes.add_textbox(Inches(1), Inches(6), Inches(11), Inches(1))
    p = contact.text_frame.paragraphs[0]
    p.text = "CoreSelf AI  |  多角色心理增效 AI 顧問  |  守住人生選擇的清醒"
    p.font.size = Pt(12)
    p.font.color.rgb = ACCENT_GOLD
    p.alignment = PP_ALIGN.CENTER

    # 儲存簡報
    prs.save('CoreSelf_AI_Presentation.pptx')
    print("簡報已成功生成：CoreSelf_AI_Presentation.pptx")

if __name__ == "__main__":
    create_premium_deck()

